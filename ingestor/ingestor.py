"""
Training material ingestion script
-----------------------------------
Reads .pptx, .docx, and .pdf files from a folder (including inside .zip
archives), chunks the text,
embeds each chunk via OpenAI, and upserts into Supabase pgvector.

Setup (run once):
  pip install python-pptx python-docx pypdf openai supabase

Environment variables:
  OPENAI_API_KEY   — from https://platform.openai.com/api-keys
  SUPABASE_URL     — from your Supabase project Settings > API
  SUPABASE_KEY     — service_role key (not anon) so we can write
"""

import io
import os
import sys
import zipfile
from pathlib import Path
from dotenv import load_dotenv
import logging
from datetime import datetime

# ── third-party ───────────────────────────────────────────────────────────────
try:
    from pptx import Presentation
    from docx import Document
    from pypdf import PdfReader
    from openai import OpenAI
    from supabase import create_client
except ImportError as e:
    sys.exit(f"Missing package: {e}\nRun: pip install python-pptx python-docx pypdf openai supabase")

SCRIPT_DIR = Path(__file__).resolve().parent
load_dotenv(SCRIPT_DIR / ".env")

logger = logging.getLogger(__name__)
logger.setLevel(logging.DEBUG)

formatter = logging.Formatter("%(name)s - %(levelname)s - %(message)s")

console_handler = logging.StreamHandler()
console_handler.setFormatter(formatter)
logger.addHandler(console_handler)

current_time = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
log_filename = SCRIPT_DIR / f"app_{current_time}.log"
file_handler = logging.FileHandler(log_filename)
file_handler.setFormatter(formatter)
logger.addHandler(file_handler)

logger.info("Application has started successfully.")

# ── config ────────────────────────────────────────────────────────────────────
DOCS_FOLDER   = "../training_docs"   # put your files here
CHUNK_WORDS   = 300                 # target words per chunk
CHUNK_OVERLAP = 50                  # words of overlap between chunks
EMBED_MODEL   = "text-embedding-3-small"

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "")
SUPABASE_URL   = os.getenv("SUPABASE_URL", "")
SUPABASE_KEY   = os.getenv("SUPABASE_KEY", "")

SUPPORTED_EXTENSIONS = {".pptx", ".docx", ".pdf"}


# ── file discovery ────────────────────────────────────────────────────────────

def _is_junk_zip_entry(name: str) -> bool:
    return (
        name.endswith("/")
        or name.startswith("__MACOSX/")
        or Path(name).name.startswith(".")
    )


def collect_documents(docs_path: Path) -> list[tuple[str, Path | io.BytesIO, str]]:
    """Return (source_label, readable_source, suffix) for each supported document."""
    documents: list[tuple[str, Path | io.BytesIO, str]] = []

    for entry in docs_path.iterdir():
        suffix = entry.suffix.lower()
        if suffix in SUPPORTED_EXTENSIONS:
            documents.append((entry.name, entry, suffix))
        elif suffix == ".zip":
            try:
                with zipfile.ZipFile(entry) as zf:
                    for name in zf.namelist():
                        if _is_junk_zip_entry(name):
                            continue
                        inner_suffix = Path(name).suffix.lower()
                        if inner_suffix not in SUPPORTED_EXTENSIONS:
                            continue
                        source = f"{entry.name}/{name}"
                        documents.append((source, io.BytesIO(zf.read(name)), inner_suffix))
            except zipfile.BadZipFile:
                logger.warning("Skipping invalid zip: %s", entry.name)

    return documents


# ── text extraction ───────────────────────────────────────────────────────────

def extract_pptx(source: Path | io.BytesIO) -> str:
    prs = Presentation(source)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def extract_docx(source: Path | io.BytesIO) -> str:
    doc = Document(source)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def extract_pdf(source: Path | io.BytesIO) -> str:
    reader = PdfReader(source)
    pages = []
    for i, page in enumerate(reader.pages, 1):
        text = (page.extract_text() or "").strip()
        if text:
            pages.append(f"[Page {i}]\n{text}")
    return "\n\n".join(pages)


def extract(source: Path | io.BytesIO, suffix: str) -> str:
    if suffix == ".pptx":
        return extract_pptx(source)
    elif suffix == ".docx":
        return extract_docx(source)
    elif suffix == ".pdf":
        return extract_pdf(source)
    return ""


# ── chunking ──────────────────────────────────────────────────────────────────

def chunk_text(text: str, chunk_words: int = CHUNK_WORDS, overlap: int = CHUNK_OVERLAP) -> list[str]:
    words = text.split()
    chunks = []
    start = 0
    while start < len(words):
        end = min(start + chunk_words, len(words))
        chunks.append(" ".join(words[start:end]))
        if end == len(words):
            break
        start += chunk_words - overlap
    return chunks


# ── embedding ─────────────────────────────────────────────────────────────────

def embed_chunks(oc: OpenAI, chunks: list[str]) -> list[list[float]]:
    """Embed chunks in batches of 96."""
    vectors = []
    batch_size = 96
    for i in range(0, len(chunks), batch_size):
        batch = chunks[i : i + batch_size]
        resp = oc.embeddings.create(model=EMBED_MODEL, input=batch)
        vectors.extend([d.embedding for d in resp.data])
    return vectors


# ── supabase upsert ───────────────────────────────────────────────────────────

def upsert(sb, source: str, chunks: list[str], vectors: list[list[float]]):
    rows = [
        {
            "source":      source,
            "chunk_index": idx,
            "content":     chunk,
            "embedding":   vector,
        }
        for idx, (chunk, vector) in enumerate(zip(chunks, vectors))
    ]
    # Delete existing rows for this source so re-running is idempotent
    sb.table("training_chunks").delete().eq("source", source).execute()
    sb.table("training_chunks").insert(rows).execute()
    logger.info("  Upserted %d chunks.", len(rows))


# ── main ──────────────────────────────────────────────────────────────────────

def main():
    if not OPENAI_API_KEY or not SUPABASE_URL or not SUPABASE_KEY:
        logger.error(
            "Missing environment variables. Set:\n"
            "  OPENAI_API_KEY\n"
            "  SUPABASE_URL\n"
            "  SUPABASE_KEY"
        )
        sys.exit(1)

    docs_path = Path(DOCS_FOLDER)
    if not docs_path.exists():
        logger.error("Folder not found: %s — create it and place your training files inside.", DOCS_FOLDER)
        sys.exit(1)

    oc = OpenAI(api_key=OPENAI_API_KEY)
    sb = create_client(SUPABASE_URL, SUPABASE_KEY)

    documents = collect_documents(docs_path)

    if not documents:
        logger.error("No supported files found in %s", DOCS_FOLDER)
        sys.exit(1)

    logger.info("Found %d file(s) to process.", len(documents))

    for source, file_source, suffix in documents:
        logger.info("Processing: %s", source)

        text = extract(file_source, suffix)
        if not text.strip():
            logger.warning("  No text extracted — skipping.")
            continue

        chunks = chunk_text(text)
        logger.info("  %d chunk(s) from %d words.", len(chunks), len(text.split()))

        vectors = embed_chunks(oc, chunks)
        upsert(sb, source, chunks, vectors)

    logger.info("Done! Your training docs are indexed and ready for the chatbot.")


if __name__ == "__main__":
    main()